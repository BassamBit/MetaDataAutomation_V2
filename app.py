import io
import re
import math
import numpy as np
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ========= Page setup =========
st.set_page_config(page_title="Excel → PowerPoint Tool", page_icon="📊", layout="centered")
st.title("Excel → PowerPoint")
st.caption("Upload your PowerPoint template and data file, then fill existing slides without duplication.")

# ========= Sidebar =========
with st.sidebar:
    st.header("Mode")
    mode = st.radio("Choose action:", ["Dictionary", "Index (TOC)"], index=0)
    st.markdown("---")

    st.header("General Settings")
    default_outname = "output.pptx"
    OUTFILE_NAME = st.text_input("Output file name", value=default_outname)

    st.markdown("---")
    st.caption("Upload template and data files below, then click Run.")

# ========= Uploads =========
pptx_file = st.file_uploader("📎 Upload PowerPoint template (.pptx)", type=["pptx"])
data_file = st.file_uploader("📎 Upload data file (Excel/CSV)", type=["xlsx", "xls", "csv"])

# Sheet picker (if Excel)
sheet_name = None
if data_file is not None and data_file.name.lower().endswith((".xlsx", ".xls")):
    try:
        xls = pd.ExcelFile(data_file)
        sheet_name = st.selectbox("Select Excel sheet", options=xls.sheet_names, index=0)
    except Exception:
        sheet_name = None

# ========= Colors (shared) =========
COLOR_DARK_BLUE = RGBColor(31, 78, 121)   # Dark Blue (Column 1 + 8 first line)
COLOR_BLACK     = RGBColor(0, 0, 0)       # Black
COLOR_GREEN     = RGBColor(82, 158, 69)   # Green (“Public” / “عام”)
COLOR_ORANGE    = RGBColor(237, 125, 49)  # Orange (“Restricted” / “مقيد”)

# ========= Helpers (shared) =========
def get_first_table(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh.table
    raise RuntimeError("No table found in this slide.")

def get_tables_sorted_by_x(slide):
    """Return all tables sorted by x-position: right-most first."""
    tables = []
    for sh in slide.shapes:
        if sh.has_table:
            tables.append((sh.left, sh.table))
    tables.sort(key=lambda t: t[0], reverse=True)  # rightmost first
    return [t[1] for t in tables]

def clear_table_data(table):
    """Clear all data rows (keep header row as-is if present)."""
    max_rows = len(table.rows) - 1
    max_cols = len(table.columns)
    for r in range(1, max_rows+1):
        for c in range(max_cols):
            table.cell(r, c).text = ""

def _clear_cell(cell):
    tf = cell.text_frame
    tf.clear()
    return tf

def _set_alignment(cell, halign=None, valign_middle=False):
    if valign_middle:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if cell.text_frame.paragraphs and halign is not None:
        cell.text_frame.paragraphs[0].alignment = halign

def _add_run(p, text, size_pt=11, color=COLOR_BLACK, bold=False):
    run = p.add_run()
    run.text = "" if text is None else str(text)
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold

def set_cell_rtl(cell, rtl=True):
    """Force RTL direction on a table cell (Arabic)."""
    txBody = cell._tc.txBody
    bodyPr = txBody.bodyPr
    bodyPr.set(qn('a:rtlCol'), '1' if rtl else '0')
    for p in cell.text_frame.paragraphs:
        pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1' if rtl else '0')

def set_cell_simple(cell, text, size=11):
    """Center both ways; black; non-bold; leave blank if empty/NaN."""
    if text is None or (isinstance(text, float) and np.isnan(text)) or str(text).strip() == "":
        cell.text = ""
        return
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = False
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def first_line_no_trailing_colon(value):
    """Take first line only; strip trailing spaces and colon variants (:, ：)."""
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return ""
    s = str(value)
    if not s.strip():
        return ""
    line = s.splitlines()[0].strip()
    line = re.sub(r'[\s:：]+$', '', line)
    return line

# ========= Dictionary formatters (unchanged) =========
def format_col1(cell, value):
    if value is None or (isinstance(value, float) and np.isnan(value)) or str(value).strip() == "":
        cell.text = ""
        return
    set_cell_rtl(cell, True)
    tf = _clear_cell(cell)
    tf.word_wrap = True
    lines = re.split(r'\r?\n', str(value))
    _set_alignment(cell, halign=PP_ALIGN.RIGHT, valign_middle=True)
    p1 = tf.paragraphs[0]
    _add_run(p1, lines[0], size_pt=12, color=COLOR_DARK_BLUE, bold=True)
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        _add_run(p, extra, size_pt=11, color=COLOR_BLACK, bold=False)
    set_cell_rtl(cell, True)

def format_col2(cell, value):
    tf = _clear_cell(cell)
    _set_alignment(cell, halign=PP_ALIGN.CENTER, valign_middle=True)
    p = tf.paragraphs[0]
    _add_run(p, value, size_pt=11, color=COLOR_BLACK, bold=False)

def format_col3(cell, value):
    text = "" if value is None else str(value).strip()
    txt_lower = text.lower()
    if text == "عام" or "public" in txt_lower:
        color = COLOR_GREEN
    elif text == "مقيد" or "restricted" in txt_lower:
        color = COLOR_ORANGE
    else:
        color = COLOR_BLACK
    tf = _clear_cell(cell)
    _set_alignment(cell, halign=PP_ALIGN.CENTER, valign_middle=True)
    p = tf.paragraphs[0]
    _add_run(p, text, size_pt=11, color=color, bold=False)

def format_col4(cell, value):
    if value is None or (isinstance(value, float) and np.isnan(value)) or str(value).strip() == "":
        cell.text = ""
        return
    tf = _clear_cell(cell)
    _set_alignment(cell, halign=PP_ALIGN.CENTER, valign_middle=True)
    p = tf.paragraphs[0]
    _add_run(p, value, size_pt=11, color=COLOR_BLACK, bold=False)

def format_col5(cell, value):
    if value is None or (isinstance(value, float) and np.isnan(value)) or str(value).strip() == "":
        cell.text = ""
        return
    tf = _clear_cell(cell)
    _set_alignment(cell, halign=PP_ALIGN.CENTER, valign_middle=True)
    p = tf.paragraphs[0]
    _add_run(p, value, size_pt=11, color=COLOR_BLACK, bold=False)

def format_col6(cell, value):
    text = "" if value is None else str(value).strip()
    txt_lower = text.lower()
    if "public" in txt_lower:
        color = COLOR_GREEN
    elif "restricted" in txt_lower:
        color = COLOR_ORANGE
    else:
        color = COLOR_BLACK
    tf = _clear_cell(cell)
    _set_alignment(cell, halign=PP_ALIGN.CENTER, valign_middle=True)
    p = tf.paragraphs[0]
    _add_run(p, text, size_pt=11, color=color, bold=False)

def format_col7(cell, value):
    tf = _clear_cell(cell)
    p = tf.paragraphs[0]
    _add_run(p, value, size_pt=11, color=COLOR_BLACK, bold=False)

def format_col8(cell, value):
    tf = _clear_cell(cell)
    lines = re.split(r'\r?\n', "" if value is None else str(value))
    _set_alignment(cell, halign=PP_ALIGN.LEFT, valign_middle=True)
    p1 = tf.paragraphs[0]
    _add_run(p1, lines[0], size_pt=12, color=COLOR_DARK_BLUE, bold=True)
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _add_run(p, extra, size_pt=11, color=COLOR_BLACK, bold=False)

def format_col9(cell, value):
    tf = _clear_cell(cell)
    p = tf.paragraphs[0]
    _add_run(p, value, size_pt=11, color=COLOR_BLACK, bold=False)

FORMATTERS_NATURAL  = {0: format_col1, 1: format_col2, 2: format_col3, 3: format_col4, 4: format_col5, 5: format_col6, 6: format_col7, 7: format_col8, 8: format_col9}
FORMATTERS_REVERSED = {8: format_col1, 7: format_col2, 6: format_col3, 5: format_col4, 4: format_col5, 3: format_col6, 2: format_col7, 1: format_col8, 0: format_col9}

# ========= Mode-specific controls =========
if mode == "Dictionary":
    st.subheader("Dictionary Filling")
    ROWS_PER_SLIDE = st.number_input("Rows per slide (Dictionary)", min_value=1, max_value=50, value=9)
    use_reversed_mapping = st.checkbox("Use reversed column mapping (8→1 ... 0→9)", value=True)

    # Column mapping UI (dropdowns) — keep original Excel order
    st.markdown("**Select columns from your Excel sheet:**")
    if data_file is not None:
        # temp read just to get columns for selects
        if data_file.name.lower().endswith(".csv"):
            tmp_df = pd.read_csv(data_file)
        else:
            tmp_df = pd.read_excel(data_file, sheet_name=sheet_name)
        cols = tmp_df.columns.tolist()

        code_col           = st.selectbox("Code column", cols)
        en_name_col        = st.selectbox("English Name column", cols)
        en_def_col         = st.selectbox("English Definition column", cols)
        classification_col = st.selectbox("Classification column", cols)
        ar_name_col        = st.selectbox("Arabic Name column", cols)
        ar_def_col         = st.selectbox("Arabic Definition column", cols)

        owner_en = st.text_input("Data Owner (English)", "Legal Department")
        owner_ar = st.text_input("Data Owner (Arabic)", "الإدارة القانونية")

else:
    st.subheader("Index (TOC) Filling")
    ROWS_PER_TABLE = st.number_input("Rows per table (TOC)", min_value=1, max_value=50, value=15)
    COLS_PER_TABLE = 4  # fixed
    st.caption("Each slide must contain two empty tables (right then left), each with a header row + data rows.")

    # TOC column selection via dropdowns
    if data_file is not None:
        if data_file.name.lower().endswith(".csv"):
            tmp_df = pd.read_csv(data_file)
        else:
            tmp_df = pd.read_excel(data_file, sheet_name=sheet_name)
        cols = tmp_df.columns.tolist()

        st.markdown("**TOC column mapping (from your DataFrame):**")
        TOC_COL1 = st.selectbox("Arabic Term+Def (first line only)", cols)   # was "المصطلح وتعريفه  "
        TOC_COL3 = st.selectbox("English Term (first line only)", cols, index=min(1, len(cols)-1))  # was "Term"
        TOC_COL4 = st.selectbox("Code (as-is)", cols, index=min(2, len(cols)-1))                    # was "Code"

# ========= Run =========
run = st.button("Run Tool ✅", type="primary", use_container_width=True)

if run:
    try:
        # 1) Load data
        if data_file is None:
            st.error("Please upload a data file (Excel/CSV).")
            st.stop()

        if data_file.name.lower().endswith(".csv"):
            df = pd.read_csv(data_file)
        else:
            df = pd.read_excel(data_file, sheet_name=sheet_name)

        st.success(f"Data loaded: {df.shape[0]} rows × {df.shape[1]} columns")

        # 2) Load template
        if pptx_file is None:
            st.error("Please upload a PowerPoint template.")
            st.stop()

        prs = Presentation(io.BytesIO(pptx_file.read()))
        st.info(f"Slides in template: {len(prs.slides)}")

        # 3) Process per mode
        if mode == "Dictionary":
            # ---- Build df2 (9 columns) in the exact order expected by FORMATTERS_NATURAL ----
            def combine_name_def(name_series, def_series, add_colon=True):
                name_series = name_series.fillna("").astype(str)
                def_series  = def_series.fillna("").astype(str)
                out = []
                for n, d in zip(name_series, def_series):
                    n = n.strip()
                    d = d.strip()
                    if not n and not d:
                        out.append("")
                    elif add_colon:
                        out.append(f"{n}:\n{d}" if d else n)
                    else:
                        out.append(f"{n}\n{d}" if d else n)
                return pd.Series(out, index=name_series.index)

            # Arabic classification derived from EN classification
            ar_class_series = df[classification_col].astype(str).replace({"Public": "عام", "Restricted": "مقيد"})

            df2 = pd.DataFrame(index=df.index)
            # Order for NATURAL mapping (0..8):
            # 0: AR Name+Def
            # 1: Owner AR
            # 2: Class AR
            # 3: Code
            # 4: Owner EN
            # 5: Class EN
            # 6: Personal Data (empty)
            # 7: EN Name+Def
            # 8: Personal Data EN (empty)
            df2[0] = df[code_col]
            df2[1] = combine_name_def(df[en_name_col], df[en_def_col], add_colon=True)
            df2[2] = owner_en
            df2[3] = df[classification_col]
            df2[4] = ""
            df2[5] = ""
            df2[6] = ar_class_series
            df2[7] = owner_ar
            df2[8] = combine_name_def(df[ar_name_col], df[ar_def_col], add_colon=True)
            # chunk rows by ROWS_PER_SLIDE
            chunks = [df2.iloc[i:i+ROWS_PER_SLIDE] for i in range(0, len(df2), ROWS_PER_SLIDE)]
            needed_slides = len(chunks)
            available_slides = len(prs.slides)

            if available_slides < needed_slides:
                st.warning(f"Template has fewer slides than needed ({available_slides} < {needed_slides}). "
                           "Only filling the available slides. Consider duplicating the template slides.")
                chunks = chunks[:available_slides]

            FORMATTERS = FORMATTERS_REVERSED if use_reversed_mapping else FORMATTERS_NATURAL

            filled_rows = 0
            for idx, chunk in enumerate(chunks):
                slide = prs.slides[idx]
                try:
                    table = get_first_table(slide)
                except RuntimeError:
                    st.warning(f"Slide {idx+1} has no table. Skipped.")
                    continue

                max_rows = len(table.rows) - 1
                max_cols = len(table.columns)
                rows_to_write = min(len(chunk), max_rows)
                cols_to_write = min(chunk.shape[1], max_cols)

                # clear table data
                clear_table_data(table)

                # write
                for r in range(rows_to_write):
                    for c in range(cols_to_write):
                        val = chunk.iat[r, c]
                        formatter = FORMATTERS.get(c, None)
                        if formatter:
                            formatter(table.cell(r+1, c), val)
                        else:
                            tf = _clear_cell(table.cell(r+1, c))
                            p = tf.paragraphs[0]
                            _add_run(p, val, size_pt=11, color=COLOR_BLACK, bold=False)
                filled_rows += rows_to_write

            if filled_rows < len(df2):
                st.warning(f"Not all rows were written ({filled_rows} / {len(df2)}). "
                           "Add more prepared slides to your template and rerun.")

        else:
            # Index (TOC)
            # Build using dropdown-selected columns
            required_cols = [TOC_COL1, TOC_COL3, TOC_COL4]
            for name in required_cols:
                if name not in df.columns:
                    st.error(f"Column not found in DataFrame: {name}")
                    st.stop()

            toc_df = df.copy()
            col1_vals = toc_df[TOC_COL1].map(first_line_no_trailing_colon)
            col2_vals = [""] * len(toc_df)  # empty on purpose
            col3_vals = toc_df[TOC_COL3].map(first_line_no_trailing_colon)
            col4_vals = toc_df[TOC_COL4]  # as-is

            # TOC matrix expected as 4 columns in this order for table:
            # (Code, EN Term, Page#, AR TermFirstLine)
            toc_matrix = list(zip(col4_vals, col3_vals, col2_vals, col1_vals))

            # Each slide contains two tables; each table has ROWS_PER_TABLE data rows
            rows_per_slide = ROWS_PER_TABLE * 2
            needed_slides = math.ceil(len(toc_matrix) / rows_per_slide)
            available_slides = len(prs.slides)

            if available_slides < needed_slides:
                st.warning(f"Template has fewer TOC slides than needed ({available_slides} < {needed_slides}). "
                           "Only filling the available slides. Consider duplicating your TOC slide.")

            total_needed_rows = len(toc_matrix)
            cur_idx = 0
            slide_counter = 0

            for slide in prs.slides:
                if slide_counter >= needed_slides:
                    break  # filled enough
                tables = get_tables_sorted_by_x(slide)
                if len(tables) < 2:
                    st.warning(f"Slide {slide_counter+1} does not contain two tables. Skipped.")
                    slide_counter += 1
                    continue

                right_table, left_table = tables[0], tables[1]
                clear_table_data(right_table)
                clear_table_data(left_table)

                # Fill right table
                for r in range(ROWS_PER_TABLE):
                    if cur_idx >= total_needed_rows:
                        break
                    row_vals = toc_matrix[cur_idx]
                    # write starting from row 0 (your template tables are headerless for data rows)
                    for c in range(COLS_PER_TABLE):
                        set_cell_simple(right_table.cell(r, c), row_vals[c], size=11)
                    cur_idx += 1
                if cur_idx >= total_needed_rows:
                    slide_counter += 1
                    break

                # Fill left table
                for r in range(ROWS_PER_TABLE):
                    if cur_idx >= total_needed_rows:
                        break
                    row_vals = toc_matrix[cur_idx]
                    for c in range(COLS_PER_TABLE):
                        set_cell_simple(left_table.cell(r, c), row_vals[c], size=11)
                    cur_idx += 1

                slide_counter += 1

            if cur_idx < total_needed_rows:
                st.warning(f"Not all TOC rows were written ({cur_idx} / {total_needed_rows}). "
                           "Add more TOC slides (two tables per slide) and rerun.")

        # 4) Save and offer download
        out_buf = io.BytesIO()
        prs.save(out_buf)
        out_buf.seek(0)
        st.success("✅ File generated successfully!")
        st.download_button(
            "⬇️ Download PowerPoint",
            data=out_buf,
            file_name=OUTFILE_NAME or "output.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

    except Exception as e:
        st.error(f"Error: {e}")
        st.stop()

st.markdown("---")
st.caption("Built with Streamlit • Supports Excel/CSV • Preserves design • Fills existing slides without creating new ones.")
